"""
main.py — UGC Monitor FastAPI 백엔드
────────────────────────────────────
Render 배포 또는 로컬 실행 모두 지원.
로컬: ../start.sh 또는 `uvicorn main:app --port 8000`
대시보드(Vercel 또는 localhost)에서 호출하면 Phase 1~3을 순서대로 실행합니다.

엔드포인트:
  POST /scan    → URL + 레퍼런스 이미지 + 프롬프트 텍스트로 전체 스캔 시작
  GET  /results → 최신 스캔 결과 조회
  GET  /health  → 서버 상태 확인
  GET  /        → 로컬 실행 시 index.html 서빙 (없으면 JSON)
"""

from __future__ import annotations  # Python 3.9 호환 (PEP 604 union syntax)

import os, io, time, json, base64, csv, requests, threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime, timezone
from typing import Optional

from fastapi import FastAPI, UploadFile, File, Form, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, FileResponse
from dotenv import load_dotenv
from PIL import Image
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
# sentence_transformers / torch 는 get_clip() 안에서 lazy import (앱 시작 속도 개선)

import gspread
from google.oauth2.service_account import Credentials

load_dotenv()

# ── 설정 ───────────────────────────────────────
APIFY_API_TOKEN    = os.getenv("APIFY_API_TOKEN")
SPREADSHEET_ID     = os.getenv("SPREADSHEET_ID")
GOOGLE_CREDENTIALS = os.getenv("GOOGLE_CREDENTIALS_PATH", "config/google_credentials.json")
NAVER_API_URL      = os.getenv("NAVER_API_URL", "").rstrip("/")
NAVER_API_KEY      = os.getenv("NAVER_API_KEY")
MY_IG_ID           = "pitapat_prompt"
SHEET_TAB_NAME     = "ugc_users"
HISTORY_TAB_NAME   = "scan_history"

APIFY_BASE     = "https://api.apify.com/v2"
ACTOR_PROFILE  = "apify~instagram-profile-scraper"
ACTOR_STORY    = "seemuapps~instagram-story-scraper"
USERNAME_HEADERS = {"username", "user", "userid", "user_id", "아이디", "id"}
MODEL_NAME     = "gemini-2.0-flash"  # Gemini 2.0 Flash via NAMC Vertex AI
API_SEMAPHORE  = threading.Semaphore(6)  # 글로벌 API 호출 동시 한도 (rate limit 회피)

PROMPT_FEED_TEMPLATE = """원본 AI 프롬프트 (이 프롬프트로 [이미지 1] 레퍼런스가 만들어짐):
═══════════════════════════════════════════════════════
{prompt_text}
═══════════════════════════════════════════════════════

[이미지 1]: 위 프롬프트로 생성된 레퍼런스 결과물
[이미지 2]: 유저가 올린 판별 대상 피드 게시물

**핵심 질문**: [이미지 2]가 위 AI 프롬프트로 다른 사람의 얼굴로 생성된 것처럼 보이나요?

판별 방법 — 위 프롬프트의 핵심 요소를 [이미지 2]에서 얼마나 만족하는지 봅니다:
1. **장면/배경**: 프롬프트가 명시한 환경과 일치하는가? (예: "차 안 뒷좌석"이면 [이미지 2]도 차 안 뒷좌석이어야 함)
2. **의상**: 프롬프트에 적힌 의상이 보이는가? (구체적인 옷 종류·색·실루엣)
3. **자세/구도**: 프롬프트가 명시한 포즈·앵글·프레이밍이 일치?
4. **색감/톤**: 프롬프트의 색감 가이드 (차가운 톤, 저채도 등) 일치?
5. **전체 인상/질감**: 프롬프트가 의도한 느낌 (예: "구형 폰카 저화질")?

**자주 오판하는 케이스 (모두 NO)**:
- 둘 다 AI풍 셀카지만 위 프롬프트의 핵심 장면이 아님 → NO
- 둘 다 자연광 인물 사진이지만 의상·구도·장소가 다름 → NO
- 위 프롬프트의 요소 중 한두 개만 부분적으로 맞음 → NO
- 단순히 "AI 셀카", "여성 인물" 같은 표면적 공통점만 겹침 → NO
- 위 프롬프트의 명시된 "절대 금지" 항목이 [이미지 2]에 보임 → NO
- 비슷해 보이지만 정말 이 프롬프트로 만든 거라 확신 안 섬 → NO

**YES 조건**: 위 프롬프트의 주요 요소들(장면, 의상, 자세, 색감, 질감) 대부분이 명확히 보이고, "이 프롬프트로 다른 사람으로 다시 생성한 결과"라고 강하게 확신될 때만.

확신 안 서면 NO. 반드시 YES 또는 NO 한 단어만 답하세요."""

PROMPT_PROFILE_TEMPLATE = """원본 AI 프롬프트 (이 프롬프트로 [이미지 1] 레퍼런스가 만들어짐):
═══════════════════════════════════════════════════════
{prompt_text}
═══════════════════════════════════════════════════════

[이미지 1]: 위 프롬프트로 생성된 고화질 레퍼런스
[이미지 2]: 유저의 프로필 사진 (저해상도 150×150, 크롭 가능)

**핵심 질문**: 이 작은 프로필 사진이 위 AI 프롬프트로 다른 사람의 얼굴로 생성된 결과물의 일부로 보이나요?

판별 방법 — 저해상도지만 다음을 확인:
1. **장면/배경의 종류**: 프롬프트가 명시한 환경(예: "차 안 뒷좌석")이 작게라도 인식되는가?
2. **의상**: 프롬프트의 의상(예: "흰 끈나시 + 회색 가디건")이 작게라도 인식되는가?
3. **포즈/앵글**: 프롬프트가 지정한 자세나 손 위치(예: "주먹으로 코+입 가림")가 보이는가?
4. **색감**: 프롬프트의 톤(예: "차가운 톤, 노란기 금지") 일치?

**저해상도이므로 매우 엄격하게**:
- 디테일이 안 보여서 확신 못 하면 무조건 NO
- 둘 다 AI풍이라는 공통점만으로는 NO
- 둘 다 여성 셀카·자연광·클로즈업 같은 표면적 공통점만으론 NO
- 배경 종류가 프롬프트와 다르면 (예: 차 안이 아니라 카페/방/거울 앞) 무조건 NO
- 위 프롬프트의 "절대 금지" 항목이 보이면 NO
- 확신이 80% 미만이면 NO

**YES 조건**: 작은 썸네일이지만 위 프롬프트의 주요 요소들이 명확히 인식되고, "이 프롬프트로 다른 사람으로 만든 결과를 작게 크롭한 것"처럼 보일 때만.

확신 안 서면 NO. 반드시 YES 또는 NO 한 단어만 답하세요."""

app = FastAPI(title="UGC Monitor API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

scan_state = {
    "status":     "idle",   # idle / running / done / error
    "progress":   0,
    "step":       "",
    "results":    [],
    "stats":      {"feed": 0, "story": 0, "profile": 0},
    "started_at": None,
    "post_url":   "",
}


def _find_file(*candidates):
    for p in candidates:
        if os.path.exists(p):
            return p
    return None


def _load_existing_phase3():
    """서버 시작 시 phase3_results.json 이 있으면 scan_state 에 복원.
    저장된 confirmed_ugc(ugc_type/feed_url) → scan_state.results(type/link/image_url) 포맷으로 변환.
    phase3_candidates.json 에서 image_url 크로스 레퍼런스."""
    results_path = _find_file("../phase3_results.json", "phase3_results.json")
    if not results_path:
        return
    try:
        with open(results_path, encoding="utf-8") as f:
            data = json.load(f)
        confirmed_raw = data.get("confirmed_ugc", [])
        if not confirmed_raw:
            return

        cand_map = {}
        cands_path = _find_file("../phase3_candidates.json", "phase3_candidates.json")
        if cands_path:
            with open(cands_path, encoding="utf-8") as f:
                for c in json.load(f):
                    cand_map[c["username"]] = c

        results = []
        for r in confirmed_raw:
            uname    = r["username"]
            ugc_type = r.get("ugc_type", "")
            link     = r.get("feed_url") or f"https://instagram.com/{uname}/"
            image_url = ""
            c = cand_map.get(uname)
            if c:
                if ugc_type == "profile":
                    image_url = c.get("profile_url", "")
                elif ugc_type == "story":
                    urls = c.get("story_image_urls") or c.get("story_urls") or []
                    image_url = urls[0] if urls else c.get("story_image_url", "")
                elif ugc_type == "feed":
                    feed_items = c.get("latest_feed_items") or c.get("feed_items") or []
                    for item in feed_items:
                        if item.get("post_url") == r.get("feed_url"):
                            image_url = item.get("image_url", "")
                            break
                    if not image_url and feed_items:
                        image_url = feed_items[0].get("image_url", "")

            results.append({
                "username":    uname,
                "detected_at": "",
                "type":        ugc_type,
                "link":        link,
                "image_url":   image_url,
                "status":      r.get("status", "pending"),  # 검토 큐용 (향후 단계)
            })

        feed_n    = sum(1 for x in results if x["type"] == "feed")
        story_n   = sum(1 for x in results if x["type"] == "story")
        profile_n = sum(1 for x in results if x["type"] == "profile")

        scan_state.update({
            "status":   "done",
            "progress": 100,
            "step":     f"기존 결과 복원 ({len(results)}명)",
            "results":  results,
            "stats":    {"feed": feed_n, "story": story_n, "profile": profile_n},
        })
        print(f"✓ phase3_results.json 복원: {len(results)}명 (feed {feed_n}, profile {profile_n}, story {story_n})")
    except Exception as e:
        print(f"⚠️  phase3_results.json 로드 실패: {e}")


_load_existing_phase3()


# ── 이미지 리사이즈 ────────────────────────────
def resize_to_data_uri(raw: bytes, max_side: int = 1024) -> str:
    """바이트 → 리사이즈된 data URI"""
    try:
        img = Image.open(io.BytesIO(raw))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        img.thumbnail((max_side, max_side), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
        return f"data:image/jpeg;base64,{b64}"
    except Exception:
        b64 = base64.b64encode(raw).decode("utf-8")
        return f"data:image/jpeg;base64,{b64}"


def video_url_to_data_uri(video_url: str) -> str | None:
    """비디오 URL → 첫 프레임 → data URI. imageio-ffmpeg 번들 바이너리 사용"""
    try:
        import subprocess
        import imageio_ffmpeg
        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        # 첫 프레임만 JPEG 파이프로 추출
        proc = subprocess.run(
            [ffmpeg, "-hide_banner", "-loglevel", "error",
             "-i", video_url, "-frames:v", "1", "-f", "image2pipe",
             "-vcodec", "mjpeg", "-"],
            capture_output=True, timeout=30,
        )
        if proc.returncode != 0 or not proc.stdout:
            return None
        img = Image.open(io.BytesIO(proc.stdout))
        if img.mode not in ("RGB", "L"):
            img = img.convert("RGB")
        img.thumbnail((1024, 1024), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=85)
        return "data:image/jpeg;base64," + base64.b64encode(buf.getvalue()).decode("utf-8")
    except Exception:
        return None


def target_to_qwen_url(url: str) -> str | None:
    """이미지면 그대로, 비디오면 첫 프레임 data URI"""
    if not url:
        return None
    lower = url.lower().split("?")[0]
    if lower.endswith((".mp4", ".mov", ".webm")):
        return video_url_to_data_uri(url)
    return url


# ── Google Sheets ─────────────────────────────
def get_sheet():
    creds = Credentials.from_service_account_file(
        GOOGLE_CREDENTIALS,
        scopes=["https://www.googleapis.com/auth/spreadsheets",
                "https://www.googleapis.com/auth/drive"],
    )
    return gspread.authorize(creds).open_by_key(SPREADSHEET_ID).worksheet(SHEET_TAB_NAME)


def sheet_range(cell: str) -> str:
    """시트명 접두사 붙인 A1 range"""
    return f"'{SHEET_TAB_NAME}'!{cell}"


# ── Apify 실행 헬퍼 ────────────────────────────
def run_apify(actor_id: str, run_input: dict, timeout: int = 200) -> list:
    resp = requests.post(
        f"{APIFY_BASE}/acts/{actor_id}/runs?token={APIFY_API_TOKEN}",
        json=run_input, timeout=30,
    )
    resp.raise_for_status()
    run_id = resp.json()["data"]["id"]

    deadline = time.time() + timeout
    status_data = {}
    while time.time() < deadline:
        time.sleep(6)
        r = requests.get(f"{APIFY_BASE}/actor-runs/{run_id}?token={APIFY_API_TOKEN}", timeout=15)
        status_data = r.json()["data"]
        if status_data["status"] == "SUCCEEDED":
            break
        if status_data["status"] in ("FAILED", "ABORTED", "TIMED-OUT"):
            return []

    dataset_id = status_data.get("defaultDatasetId", "")
    if not dataset_id:
        return []

    return requests.get(
        f"{APIFY_BASE}/datasets/{dataset_id}/items?token={APIFY_API_TOKEN}",
        timeout=30,
    ).json()


# ── NAVER Open Models (Qwen2.5-VL) 판별 ───────
def call_model_single(reference_data_uri: str, target_url: str, img_type: str = "feed",
                      prompt_text: str = "", max_retries: int = 5) -> bool | None:
    """1개 레퍼런스 vs 1개 타겟 비교 (Gemini 2.0 Flash via NAMC Vertex AI).
    prompt_text가 있으면 원본 AI 프롬프트를 함께 전달 (하이브리드 판별)."""
    template = PROMPT_PROFILE_TEMPLATE if img_type in ("profile", "story") else PROMPT_FEED_TEMPLATE
    prompt = template.format(prompt_text=prompt_text or "(프롬프트 텍스트 없음 — 이미지 비교만 수행)")
    target = target_to_qwen_url(target_url)
    if not target:
        return None
    payload = {
        "model": MODEL_NAME,
        "target_model_names": MODEL_NAME,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "text", "text": prompt},
                {"type": "text", "text": "[이미지 1] 레퍼런스:"},
                {"type": "image_url", "image_url": {"url": reference_data_uri}},
                {"type": "text", "text": "[이미지 2] 판별 대상:"},
                {"type": "image_url", "image_url": {"url": target}},
            ],
        }],
        "temperature": 0.1,
        "max_tokens": 10,
    }
    headers = {
        "Authorization": f"Bearer {NAVER_API_KEY}",
        "custom-llm-provider": "vertex_ai",
        "Content-Type": "application/json",
    }
    endpoint = f"{NAVER_API_URL}/chat/completions"
    for attempt in range(max_retries):
        try:
            with API_SEMAPHORE:
                resp = requests.post(endpoint, json=payload, headers=headers, timeout=90)
            if resp.status_code in (429, 503):
                time.sleep(min(60, 2 ** (attempt + 2)))
                continue
            resp.raise_for_status()
            answer = resp.json()["choices"][0]["message"]["content"].strip().upper()
            return "YES" in answer
        except Exception as e:
            print(f"⚠️ Gemini 호출 실패: {str(e)[:120]}")
            return None
    return None


def call_qwen(reference_data_uris: list, target_url: str, img_type: str = "feed",
              prompt_text: str = "") -> bool:
    """다수의 레퍼런스 vs 1개 타겟 — 각 ref마다 별도 호출 후 2/3 다수결로 매치 결정.
    이름은 후방 호환을 위해 call_qwen 유지하지만 실제로는 Gemini 2.0 Flash 사용.
    레퍼런스 1장이면 단순 매치, 2장 이상이면 절반 이상이 YES일 때 매치."""
    if not reference_data_uris:
        return False
    with ThreadPoolExecutor(max_workers=len(reference_data_uris)) as ex:
        futs = [ex.submit(call_model_single, ru, target_url, img_type, prompt_text)
                for ru in reference_data_uris]
        results = [f.result() for f in as_completed(futs)]
    yes_count = sum(1 for r in results if r is True)
    threshold = max(1, (len(reference_data_uris) + 1) // 2)  # 다수결 (3장이면 2개, 1장이면 1개)
    return yes_count >= threshold


# ── CLIP 기반 이미지 판별 ───────────────────────
# 타입별 threshold: 피드는 엄격, 프사/스토리는 관대 (저해상도/분할 감안)
CLIP_THRESHOLDS = {"feed": 0.78, "profile": 0.72, "story": 0.68}

_clip_model = None
_clip_util  = None
_clip_lock  = threading.Lock()

def get_clip():
    """Lazy-load CLIP 모델 + sentence_transformers 자체도 여기서 import"""
    global _clip_model, _clip_util
    if _clip_model is None:
        with _clip_lock:
            if _clip_model is None:
                print("🔧 sentence_transformers import 중...")
                t0 = time.time()
                from sentence_transformers import SentenceTransformer, util
                print(f"   import 완료 ({time.time()-t0:.1f}s)")
                print("🔧 CLIP 모델 로드 중...")
                t1 = time.time()
                _clip_model = SentenceTransformer("clip-ViT-B-32")
                _clip_util  = util
                print(f"   로드 완료 ({time.time()-t1:.1f}s)")
    return _clip_model


def load_image_from_data_uri(uri: str) -> Image.Image | None:
    try:
        b64 = uri.split(",", 1)[1] if "," in uri else uri
        img = Image.open(io.BytesIO(base64.b64decode(b64)))
        return img.convert("RGB") if img.mode != "RGB" else img
    except Exception as e:
        print(f"⚠️ data URI 디코딩 실패: {e}")
        return None


def fetch_image_for_clip(url: str) -> Image.Image | None:
    """URL에서 이미지 다운로드 → PIL Image. 비디오면 첫 프레임."""
    if not url:
        return None
    lower = url.lower().split("?")[0]
    if lower.endswith((".mp4", ".mov", ".webm")):
        data_uri = video_url_to_data_uri(url)
        return load_image_from_data_uri(data_uri) if data_uri else None
    try:
        r = requests.get(url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
        if r.status_code == 200:
            img = Image.open(io.BytesIO(r.content))
            return img.convert("RGB") if img.mode != "RGB" else img
    except Exception:
        pass
    return None


def compute_ref_embeddings(reference_data_uris: list):
    """레퍼런스 이미지들의 CLIP 임베딩 (스캔 시작 시 1회)"""
    model = get_clip()
    imgs  = [load_image_from_data_uri(u) for u in reference_data_uris]
    imgs  = [i for i in imgs if i is not None]
    if not imgs:
        return None
    return model.encode(imgs, convert_to_tensor=True, show_progress_bar=False)


def clip_matches(ref_embs, target_url: str, img_type: str) -> bool:
    """타겟 이미지가 레퍼런스 중 하나라도 임계값 이상 유사하면 True"""
    target_img = fetch_image_for_clip(target_url)
    if target_img is None:
        return False
    try:
        model      = get_clip()
        target_emb = model.encode(target_img, convert_to_tensor=True, show_progress_bar=False)
        max_sim    = _clip_util.cos_sim(target_emb, ref_embs).max().item()
        threshold  = CLIP_THRESHOLDS.get(img_type, 0.75)
        return max_sim >= threshold
    except Exception as e:
        print(f"⚠️ CLIP 비교 실패: {e}")
        return False


# ── 스캔 히스토리 저장 ─────────────────────────
def save_scan_history(post_url: str, stats: dict, confirmed: list,
                      campaign_name: str = "", reviewer: str = ""):
    try:
        creds = Credentials.from_service_account_file(
            GOOGLE_CREDENTIALS,
            scopes=["https://www.googleapis.com/auth/spreadsheets",
                    "https://www.googleapis.com/auth/drive"],
        )
        ss = gspread.authorize(creds).open_by_key(SPREADSHEET_ID)
        try:
            ws = ss.worksheet(HISTORY_TAB_NAME)
            header = ws.row_values(1)
            if "캠페인" not in header:
                ws.update_cell(1, 2, "캠페인")
            if "실행자" not in header:
                # 새 컬럼 추가 (9번째 자리)
                ws.update_cell(1, 9, "실행자")
        except gspread.WorksheetNotFound:
            ws = ss.add_worksheet(HISTORY_TAB_NAME, rows=1000, cols=9)
            ws.append_row(["날짜", "캠페인", "게시물URL", "피드", "스토리", "프사", "총계", "유저목록", "실행자"],
                          value_input_option="RAW")
        now = datetime.now().strftime("%Y-%m-%d %H:%M")
        usernames = ",".join(r["username"] for r in confirmed)
        ws.append_row([
            now, campaign_name, post_url,
            stats.get("feed", 0), stats.get("story", 0), stats.get("profile", 0),
            len(confirmed), usernames, reviewer,
        ], value_input_option="RAW")
    except Exception as e:
        print(f"⚠️ 히스토리 저장 실패: {e}")


# ── 전체 스캔 파이프라인 ───────────────────────
def parse_comment_file(content: bytes, filename: str) -> list[str]:
    """xlsx 또는 csv에서 username 리스트 추출"""
    ext = os.path.splitext(filename)[1].lower()
    if ext == ".xlsx":
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        ws = wb.active
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            return []
        header = [str(c).strip().lower() if c else "" for c in rows[0]]
        col_idx = next((i for i, h in enumerate(header) if h in USERNAME_HEADERS), 0)
        return [str(r[col_idx]).strip() for r in rows[1:] if r and r[col_idx]]
    elif ext == ".csv":
        text = content.decode("utf-8-sig")
        reader = list(csv.reader(io.StringIO(text)))
        if not reader:
            return []
        header = [c.strip().lower() for c in reader[0]]
        col_idx = next((i for i, h in enumerate(header) if h in USERNAME_HEADERS), 0)
        return [r[col_idx].strip() for r in reader[1:] if r and r[col_idx]]
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")


def run_full_scan(comment_file_bytes: bytes, comment_filename: str,
                  reference_data_uris: list, post_url: str = "",
                  prompt_text: str = "", campaign_name: str = "",
                  reviewer: str = ""):
    global scan_state
    scan_state.update({
        "status": "running", "progress": 5, "step": "댓글 파일 파싱 중...",
        "results": [], "started_at": datetime.now().isoformat(),
        "campaign_name": campaign_name,
        "reviewer": reviewer,
    })

    try:
        sheet = get_sheet()

        # ── Phase 1: 댓글 파일 파싱 & 시트 추가 ──
        raw_usernames = parse_comment_file(comment_file_bytes, comment_filename)
        # 본인 제외 + 중복 제거
        seen = set()
        usernames = []
        for u in raw_usernames:
            if not u or u.lower() == MY_IG_ID.lower() or u.lower() in seen:
                continue
            seen.add(u.lower())
            usernames.append(u)

        scan_state.update({"progress": 15, "step": f"댓글 유저 {len(usernames)}명 확인"})

        existing = {r[0].strip() for r in sheet.get_all_values()[1:] if r and r[0]}
        to_add   = [u for u in usernames if u not in existing]
        if to_add:
            now  = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
            rows = [[u,"",now,"","","","","none","",post_url,""] for u in to_add]
            sheet.append_rows(rows, value_input_option="RAW")

        scan_state.update({"progress": 25, "step": f"프로필 스캔 중... ({len(usernames)}명)"})

        # ── Phase 2: 프로필 스캔 ──────────────
        profiles  = []
        chunks    = [usernames[i:i+50] for i in range(0, len(usernames), 50)]

        for idx, chunk in enumerate(chunks):
            p = run_apify(ACTOR_PROFILE, {
                "usernames": chunk, "resultsLimit": 1,
                "_triggeredBy": "지원", "_project": "프롬프트 오가닉 모니터링",
            })
            profiles.extend(p)
            progress = 30 + int((idx+1) / len(chunks) * 30)
            scan_state.update({
                "progress": progress,
                "step": f"프로필 스캔 중... ({(idx+1)*50}/{len(usernames)}명)",
            })

        profile_map = {}
        for p in profiles:
            uname = p.get("username", "")
            if uname:
                profile_map[uname.lower()] = p

        # ── Phase 2b: 스토리 스캔 (hasPublicStory 유저만) ──
        story_users = [u for u in usernames
                       if profile_map.get(u.lower(), {}).get("hasPublicStory", False)]
        scan_state.update({
            "progress": 62,
            "step": f"스토리 스캔 중... ({len(story_users)}명 활성 스토리)",
        })
        story_map = {}
        story_chunks = [story_users[i:i+20] for i in range(0, len(story_users), 20)]
        for idx, chunk in enumerate(story_chunks):
            try:
                items = run_apify(ACTOR_STORY, {"usernames": chunk}, timeout=180)
                for it in items:
                    u = (it.get("username") or "").lower()
                    stories = it.get("stories") or []
                    urls = [s.get("mediaUrl") for s in stories if s.get("mediaUrl")]
                    if u and urls:
                        story_map[u] = urls
            except Exception as e:
                print(f"story batch {idx+1} 실패: {e}")
            if idx < len(story_chunks) - 1:
                time.sleep(2)

        # 판별 후보 구성 (피드 2장, 스토리 2장으로 제한)
        candidates = []
        for uname, p in profile_map.items():
            story_urls = story_map.get(uname.lower(), [])         # 전체
            latest_posts = p.get("latestPosts") or p.get("posts") or []
            feed_items   = []
            for lp in latest_posts[:3]:                          # max 3
                img_url = lp.get("displayUrl") or lp.get("imageUrl") or ""
                sc      = lp.get("shortCode") or lp.get("shortcode") or ""
                p_url   = f"https://www.instagram.com/p/{sc}/" if sc else lp.get("url", "")
                if img_url:
                    feed_items.append({"image_url": img_url, "post_url": p_url})
            profile_url = p.get("profilePicUrl") or p.get("profilePicUrlHD", "")
            if story_urls or feed_items or profile_url:
                candidates.append({
                    "username":    uname,
                    "story_urls":  story_urls,
                    "feed_items":  feed_items,
                    "profile_url": profile_url,
                })

        scan_state.update({"progress": 65, "step": f"AI 이미지 판별 중... (0/{len(candidates)}명)"})

        # ── Phase 3: NAVER Qwen3.5-35B-A3B (MoE) 판별 ───
        confirmed     = []
        done_count    = 0
        done_lock     = threading.Lock()

        def detect_one(user):
            """단일 유저 판별 — NAVER API 호출"""
            images = []
            if user.get("profile_url"):
                images.append(("profile", user["profile_url"], ""))
            for s_url in user.get("story_urls", []):
                images.append(("story", s_url, ""))
            for item in user.get("feed_items", []):
                images.append(("feed", item["image_url"], item.get("post_url", "")))

            for img_type, img_url, p_url in images:
                if call_qwen(reference_data_uris, img_url, img_type, prompt_text) is True:
                    return {
                        "username":    user["username"],
                        "detected_at": datetime.now().strftime("%H:%M"),
                        "type":        img_type,
                        "link":        p_url or None,
                        "image_url":   img_url,
                    }
            return None

        with ThreadPoolExecutor(max_workers=8) as pool:
            futures = {pool.submit(detect_one, u): u for u in candidates}
            for future in as_completed(futures):
                with done_lock:
                    done_count += 1
                    n = done_count
                result = future.result()
                if result:
                    confirmed.append(result)
                scan_state.update({
                    "progress": 65 + int(n / max(len(candidates), 1) * 30),
                    "step": f"AI 이미지 판별 중... ({n}/{len(candidates)}명)",
                })

        # Sheets 일괄 업데이트 (완료 후 한 번에)
        if confirmed:
            all_vals  = sheet.get_all_values()
            row_index = {row[0].strip().lower(): i+2
                         for i, row in enumerate(all_vals[1:]) if row and row[0]}
            now_str   = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
            batch     = []
            for r in confirmed:
                ri = row_index.get(r["username"].lower())
                if ri:
                    batch += [
                        {"range": sheet_range(f"G{ri}"), "values": [["TRUE"]]},
                        {"range": sheet_range(f"H{ri}"), "values": [[r["type"]]]},
                        {"range": sheet_range(f"I{ri}"), "values": [[now_str]]},
                    ]
            if batch:
                sheet.spreadsheet.values_batch_update(
                    {"valueInputOption": "RAW", "data": batch}
                )

        # ── 완료 ──────────────────────────────
        feed_n    = len([r for r in confirmed if r["type"] == "feed"])
        story_n   = len([r for r in confirmed if r["type"] == "story"])
        profile_n = len([r for r in confirmed if r["type"] == "profile"])
        stats_final = {"feed": feed_n, "story": story_n, "profile": profile_n}

        save_scan_history(post_url, stats_final, confirmed, campaign_name, reviewer)

        scan_state.update({
            "status":   "done",
            "progress": 100,
            "step":     "완료!",
            "results":  confirmed,
            "stats":    stats_final,
        })

    except Exception as e:
        scan_state.update({"status": "error", "step": f"오류: {str(e)}"})


# ── API 엔드포인트 ─────────────────────────────
@app.get("/health")
def health():
    return {"status": "ok", "time": datetime.now().isoformat()}


@app.post("/scan")
async def start_scan(
    background_tasks: BackgroundTasks,
    comment_file: UploadFile = File(...),
    post_url: str = Form(""),
    prompt_text: str = Form(""),
    campaign_name: str = Form(""),
    reviewer: str = Form(""),
    reference_image_1: Optional[UploadFile] = File(None),
    reference_image_2: Optional[UploadFile] = File(None),
    reference_image_3: Optional[UploadFile] = File(None),
    reference_image_4: Optional[UploadFile] = File(None),
    reference_image_5: Optional[UploadFile] = File(None),
):
    if scan_state["status"] == "running":
        return JSONResponse({"error": "이미 스캔이 진행 중입니다."}, status_code=409)

    ref_files = [f for f in [reference_image_1, reference_image_2, reference_image_3,
                              reference_image_4, reference_image_5] if f is not None]
    if not ref_files:
        return JSONResponse({"error": "레퍼런스 이미지를 1장 이상 업로드해주세요."}, status_code=422)

    if not prompt_text or not prompt_text.strip():
        return JSONResponse({"error": "원본 AI 프롬프트 텍스트를 입력해주세요."}, status_code=422)

    ref_uris = []
    for f in ref_files:
        raw = await f.read()
        ref_uris.append(resize_to_data_uri(raw))

    comment_bytes = await comment_file.read()
    filename      = comment_file.filename or ""

    # 백그라운드 작업 등록 전에 즉시 running으로 리셋 — 폴링이 먼저 돌 때 이전 done 상태를 읽지 않도록
    scan_state.update({
        "status": "running", "progress": 1, "step": "시작 중...",
        "results": [], "stats": {"feed": 0, "story": 0, "profile": 0},
        "started_at": datetime.now().isoformat(),
        "post_url": post_url,
        "campaign_name": campaign_name,
        "reviewer": reviewer,
    })

    background_tasks.add_task(run_full_scan, comment_bytes, filename, ref_uris,
                              post_url, prompt_text, campaign_name, reviewer)
    return {"status": "started", "filename": filename, "ref_count": len(ref_uris),
            "post_url": post_url, "prompt_chars": len(prompt_text),
            "campaign_name": campaign_name, "reviewer": reviewer}


@app.get("/results")
def get_results():
    return scan_state


REVIEW_LOG_TAB = "review_log"
REVIEW_LOG_HEADER = ["timestamp", "username", "type", "decision", "reviewer", "campaign"]


def _get_reviewer() -> str:
    """검수자 식별 — .env 의 REVIEWER_EMAIL 우선, 없으면 macOS 사용자명"""
    return os.getenv("REVIEWER_EMAIL") or os.getenv("USER") or "unknown"


def _log_review_to_sheets(username: str, ugc_type: str, decision: str, reviewer: str = "") -> None:
    """review_log 탭에 검수 이력 1행 append. 탭 없으면 헤더와 함께 생성."""
    try:
        creds = Credentials.from_service_account_file(
            GOOGLE_CREDENTIALS,
            scopes=["https://www.googleapis.com/auth/spreadsheets",
                    "https://www.googleapis.com/auth/drive"],
        )
        ss = gspread.authorize(creds).open_by_key(SPREADSHEET_ID)
        try:
            ws = ss.worksheet(REVIEW_LOG_TAB)
        except gspread.WorksheetNotFound:
            ws = ss.add_worksheet(title=REVIEW_LOG_TAB, rows=1000, cols=len(REVIEW_LOG_HEADER))
            ws.append_row(REVIEW_LOG_HEADER, value_input_option="RAW")
        ws.append_row([
            datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
            username,
            ugc_type,
            decision,
            reviewer or scan_state.get("reviewer") or _get_reviewer(),
            scan_state.get("campaign_name", "") or os.getenv("CAMPAIGN_NAME", ""),
        ], value_input_option="RAW")
    except Exception as e:
        print(f"⚠️  review_log Sheets 기록 실패: {e}")


def _save_phase3_results() -> None:
    """scan_state.results 의 status 변경을 phase3_results.json 에 반영."""
    path = _find_file("../phase3_results.json", "phase3_results.json")
    if not path:
        return
    try:
        with open(path, encoding="utf-8") as f:
            data = json.load(f)
    except Exception:
        data = {"confirmed_ugc": [], "all_results": []}
    by_user = {r["username"]: r.get("status", "pending") for r in scan_state["results"]}
    for r in data.get("confirmed_ugc", []):
        if r["username"] in by_user:
            r["status"] = by_user[r["username"]]
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"⚠️  phase3_results.json 저장 실패: {e}")


@app.post("/review/decide")
def review_decide(username: str = Form(...), decision: str = Form(...),
                  reviewer: str = Form("")):
    if decision not in ("approved", "rejected"):
        return JSONResponse({"error": "decision must be 'approved' or 'rejected'"}, status_code=422)
    hit = None
    for r in scan_state.get("results", []):
        if r.get("username") == username:
            r["status"] = decision
            hit = r
            break
    if not hit:
        return JSONResponse({"error": f"username '{username}' not found in current results"}, status_code=404)
    if reviewer:
        scan_state["reviewer"] = reviewer  # 브라우저에서 보낸 이름을 메모리에 저장
    _save_phase3_results()
    _log_review_to_sheets(username, hit.get("type", ""), decision, reviewer)
    return {"ok": True, "username": username, "decision": decision, "reviewer": reviewer}


@app.get("/review/pending")
def review_pending():
    """검토 대기(status=pending) 항목만 필터링해서 반환."""
    results = scan_state.get("results", [])
    return {
        "pending":  [r for r in results if r.get("status") == "pending"],
        "approved": [r for r in results if r.get("status") == "approved"],
        "rejected": [r for r in results if r.get("status") == "rejected"],
    }


PHASE3_MATCHED_TAB = "phase3_matched"
PHASE3_MATCHED_HEADER = ["timestamp", "username", "type", "link", "campaign", "reviewer"]


@app.post("/review/export")
def review_export():
    """승인된(approved) 유저들을 Google Sheets phase3_matched 탭에 일괄 기록.
    이미 있는 username 은 중복 기록하지 않음."""
    results = scan_state.get("results", [])
    approved = [r for r in results if r.get("status") == "approved"]
    if not approved:
        return JSONResponse({"error": "승인된 유저 없음"}, status_code=422)
    try:
        creds = Credentials.from_service_account_file(
            GOOGLE_CREDENTIALS,
            scopes=["https://www.googleapis.com/auth/spreadsheets",
                    "https://www.googleapis.com/auth/drive"],
        )
        ss = gspread.authorize(creds).open_by_key(SPREADSHEET_ID)
        try:
            ws = ss.worksheet(PHASE3_MATCHED_TAB)
        except gspread.WorksheetNotFound:
            ws = ss.add_worksheet(title=PHASE3_MATCHED_TAB, rows=500, cols=len(PHASE3_MATCHED_HEADER))
            ws.append_row(PHASE3_MATCHED_HEADER, value_input_option="RAW")
        existing = {row[1] for row in ws.get_all_values()[1:] if len(row) > 1 and row[1]}
        now_str = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
        reviewer = scan_state.get("reviewer", "") or _get_reviewer()
        campaign = scan_state.get("campaign_name", "")
        new_rows = []
        for r in approved:
            if r["username"] in existing:
                continue
            new_rows.append([now_str, r["username"], r.get("type", ""),
                             r.get("link", ""), campaign, reviewer])
        if new_rows:
            ws.append_rows(new_rows, value_input_option="RAW")
        return {"exported": len(new_rows), "skipped_duplicates": len(approved) - len(new_rows),
                "tab": PHASE3_MATCHED_TAB}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/history")
def get_history():
    try:
        creds = Credentials.from_service_account_file(
            GOOGLE_CREDENTIALS,
            scopes=["https://www.googleapis.com/auth/spreadsheets",
                    "https://www.googleapis.com/auth/drive"],
        )
        ss = gspread.authorize(creds).open_by_key(SPREADSHEET_ID)
        try:
            ws = ss.worksheet(HISTORY_TAB_NAME)
        except gspread.WorksheetNotFound:
            return {"history": []}
        rows = ws.get_all_values()
        if len(rows) <= 1:
            return {"history": []}
        keys = rows[0]
        return {"history": [dict(zip(keys, r)) for r in rows[1:]]}
    except Exception as e:
        return {"history": [], "error": str(e)}


def _build_pptx(results: list, stats: dict, date_str: str, post_url: str) -> io.BytesIO:
    """스캔 결과 → PPTX BytesIO"""
    TEAL   = RGBColor(0x5B, 0xBF, 0xAD); TEAL_D = RGBColor(0x3D, 0x9E, 0x8E)
    AMBER  = RGBColor(0xE0, 0x9A, 0x5A); AMBER_D= RGBColor(0xC0, 0x78, 0x40)
    BLUE   = RGBColor(0x6A, 0x9F, 0xD8); BLUE_D = RGBColor(0x48, 0x78, 0xB8)
    DARK   = RGBColor(0x1C, 0x1C, 0x1A); MED    = RGBColor(0x4A, 0x4A, 0x48)
    GRAY   = RGBColor(0x8A, 0x88, 0x80); LGRAY  = RGBColor(0xC4, 0xC2, 0xBA)
    WHITE  = RGBColor(0xFF, 0xFF, 0xFF); BG     = RGBColor(0xF7, 0xF6, 0xF2)
    SURF   = RGBColor(0xFF, 0xFF, 0xFF); BORDER = RGBColor(0xE2, 0xE0, 0xD8)

    TYPE_CFG = {
        "feed":    {"ko": "피드",   "c": TEAL,  "d": TEAL_D,  "bg": RGBColor(0xE4,0xF5,0xF2)},
        "story":   {"ko": "스토리", "c": AMBER, "d": AMBER_D, "bg": RGBColor(0xFD,0xF0,0xE2)},
        "profile": {"ko": "프사",   "c": BLUE,  "d": BLUE_D,  "bg": RGBColor(0xE6,0xF0,0xFB)},
    }

    def txt(slide, text, l, t, w, h, size, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, italic=False):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = text; run.font.name = "Arial"
        run.font.size = Pt(size); run.font.bold = bold
        run.font.italic = italic; run.font.color.rgb = color

    def rct(slide, l, t, w, h, fill, line=None):
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = fill
        if line: s.line.color.rgb = line
        else:    s.line.fill.background()
        return s

    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Cover ────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    s1.background.fill.solid(); s1.background.fill.fore_color.rgb = BG
    rct(s1, 0, 0, prs.slide_width, Inches(0.08), TEAL)
    rct(s1, Inches(1.0), Inches(0.5), Inches(0.06), Inches(1.5), TEAL)
    txt(s1, "UGC MONITORING REPORT",
        Inches(1.2), Inches(0.55), Inches(10), Inches(0.4), size=10, color=GRAY)
    txt(s1, "pitapat_prompt",
        Inches(1.2), Inches(0.95), Inches(10), Inches(0.9), size=36, bold=True, color=DARK)
    txt(s1, date_str, Inches(1.2), Inches(1.9), Inches(5), Inches(0.4), size=13, color=GRAY)
    if post_url:
        txt(s1, post_url, Inches(1.2), Inches(2.3), Inches(9), Inches(0.4), size=12, color=TEAL_D)
    rct(s1, Inches(1.0), Inches(2.85), Inches(11.3), Inches(0.015), BORDER)

    card_data = [
        ("피드 UGC",   stats.get("feed",    0), "feed"),
        ("스토리 UGC", stats.get("story",   0), "story"),
        ("프사 변경",  stats.get("profile", 0), "profile"),
    ]
    for (label, count, ttype), left in zip(card_data, [Inches(1.0), Inches(4.6), Inches(8.2)]):
        cfg = TYPE_CFG[ttype]
        rct(s1, left, Inches(3.1), Inches(3.3), Inches(2.5), SURF, BORDER)
        rct(s1, left, Inches(3.1), Inches(3.3), Inches(0.07), cfg["c"])
        txt(s1, label, left+Inches(0.22), Inches(3.32), Inches(2.9), Inches(0.4), size=11, color=GRAY)
        txt(s1, str(count), left+Inches(0.18), Inches(3.72), Inches(2.9), Inches(1.2),
            size=60, bold=True, color=cfg["d"])
        txt(s1, "건", left+Inches(0.22), Inches(4.95), Inches(2.9), Inches(0.4), size=13, color=GRAY)
    txt(s1, f"총  {len(results)}건  감지",
        Inches(1.0), Inches(5.85), Inches(11.3), Inches(0.55),
        size=15, bold=True, color=MED, align=PP_ALIGN.CENTER)

    # ── Slides 2+: Per-UGC ────────────────────────────────────
    for r in results:
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid(); sl.background.fill.fore_color.rgb = BG
        cfg   = TYPE_CFG.get(r.get("type", "feed"), TYPE_CFG["feed"])
        uname = r.get("username", "")
        link  = r.get("link") or ""
        time_ = r.get("detected_at", "")

        rct(sl, 0, 0, prs.slide_width, Inches(0.07), cfg["c"])

        # 이미지
        img_ok = False
        img_url = r.get("image_url", "")
        if img_url:
            try:
                ir = requests.get(img_url, timeout=10, headers={"User-Agent": "Mozilla/5.0"})
                if ir.status_code == 200:
                    rct(sl, Inches(0.56), Inches(0.41), Inches(7.4), Inches(6.9), LGRAY)
                    sl.shapes.add_picture(io.BytesIO(ir.content),
                                          Inches(0.5), Inches(0.35), Inches(7.4), Inches(6.9))
                    img_ok = True
            except Exception:
                pass
        if not img_ok:
            rct(sl, Inches(0.5), Inches(0.35), Inches(7.4), Inches(6.9), BORDER)
            txt(sl, "이미지 없음", Inches(0.5), Inches(3.5), Inches(7.4), Inches(0.5),
                size=14, color=LGRAY, align=PP_ALIGN.CENTER)

        PNL = Inches(8.15); PW = Inches(4.7)
        rct(sl, PNL, Inches(0.45), Inches(1.3), Inches(0.42), cfg["bg"], cfg["c"])
        txt(sl, cfg["ko"], PNL, Inches(0.46), Inches(1.3), Inches(0.40),
            size=12, bold=True, color=cfg["d"], align=PP_ALIGN.CENTER)
        txt(sl, f"@{uname}", PNL, Inches(1.1), PW, Inches(0.75), size=26, bold=True, color=DARK)
        txt(sl, f"감지  {time_}", PNL, Inches(1.95), PW, Inches(0.4), size=12, color=GRAY)
        rct(sl, PNL, Inches(2.55), PW, Inches(0.018), BORDER)
        if link:
            txt(sl, "게시물 링크", PNL, Inches(2.75), PW, Inches(0.35), size=10, color=LGRAY)
            txt(sl, link, PNL, Inches(3.1), PW, Inches(0.55), size=11, color=cfg["d"])
        else:
            txt(sl, "링크 없음 (스토리/프사)", PNL, Inches(2.75), PW, Inches(0.4),
                size=11, color=LGRAY, italic=True)
        rct(sl, PNL, Inches(6.9), PW, Inches(0.018), BORDER)
        txt(sl, {"feed":"피드 게시물","story":"스토리","profile":"프로필 사진"}.get(r.get("type","feed"),""),
            PNL, Inches(7.0), PW, Inches(0.35), size=10, color=LGRAY)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


@app.get("/export/slides")
def export_slides():
    results  = scan_state.get("results", [])
    stats    = scan_state.get("stats",   {"feed": 0, "story": 0, "profile": 0})
    started  = scan_state.get("started_at", "")
    post_url = scan_state.get("post_url", "")
    date_str = started[:10] if started else datetime.now().strftime("%Y-%m-%d")

    pptx_buf = _build_pptx(results, stats, date_str, post_url)

    try:
        creds = Credentials.from_service_account_file(
            GOOGLE_CREDENTIALS,
            scopes=["https://www.googleapis.com/auth/drive"],
        )
        drive = build("drive", "v3", credentials=creds)

        file_meta = {
            "name": f"UGC 리포트 {date_str}",
            "mimeType": "application/vnd.google-apps.presentation",
        }
        media = MediaIoBaseUpload(
            pptx_buf,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        file = drive.files().create(
            body=file_meta, media_body=media, fields="id,webViewLink"
        ).execute()

        drive.permissions().create(
            fileId=file["id"],
            body={"type": "anyone", "role": "writer"},
        ).execute()

        return {"url": file["webViewLink"]}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


_INDEX_HTML_PATH = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "index.html",
)

@app.get("/")
def root():
    # 로컬 실행 시 같은 폴더 구조면 index.html 서빙 (대시보드 즉시 사용 가능)
    if os.path.exists(_INDEX_HTML_PATH):
        return FileResponse(_INDEX_HTML_PATH)
    return {"message": "UGC Monitor API", "docs": "/docs"}
